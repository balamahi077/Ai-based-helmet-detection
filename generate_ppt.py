#!/usr/bin/env python3
"""
Generate a PPT for the AI Safety Helmet Detection Project.
Output: AI_Helmet_Detection_Project.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
import os

PROJECT_TITLE = "AI-based Safety Helmet Detection"
OUTPUT_FILE = "AI_Helmet_Detection_Project.pptx"


def add_title_slide(prs):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = PROJECT_TITLE
    slide.placeholders[1].text = "Computer Vision • Flask • YOLOv8 • OpenCV"


def add_bulleted_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = bullet
        p.level = 0
    return slide


def add_content_slide(prs, title, left_content, right_content=None):
    slide_layout = prs.slide_layouts[3]  # Two Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    left = slide.placeholders[1].text_frame
    left.clear()
    for i, line in enumerate(left_content):
        p = left.add_paragraph() if i > 0 else left.paragraphs[0]
        p.text = line
    if right_content is not None:
        right = slide.placeholders[2].text_frame
        right.clear()
        for i, line in enumerate(right_content):
            p = right.add_paragraph() if i > 0 else right.paragraphs[0]
            p.text = line
    return slide


def add_image_slide(prs, title, image_path):
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(1), Inches(1.5), width=Inches(8))
    else:
        body = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
        tf = body.text_frame
        tf.text = f"Image not found: {image_path}"
    return slide


def build_presentation():
    prs = Presentation()

    # Title
    add_title_slide(prs)

    # Problem & Solution
    add_bulleted_slide(
        prs,
        "Problem & Solution",
        [
            "Problem: Ensure workplace safety by verifying helmet usage",
            "Solution: AI-powered detection for images, videos, and live streams",
            "Benefits: Real-time alerts, analytics, compliance tracking",
        ],
    )

    # Features
    add_bulleted_slide(
        prs,
        "Key Features",
        [
            "Image detection with annotations",
            "Video analysis and compliance metrics",
            "Real-time camera detection",
            "Modern dashboard (HTML/CSS/JS, no React)",
            "Exportable statistics & logs",
        ],
    )

    # Architecture
    add_content_slide(
        prs,
        "Architecture",
        [
            "Frontend: templates/index.html, templates/realtime.html",
            "Backend: Flask app (app.py)",
            "Detection Engine: helmet_detector.py",
            "Models: YOLOv8 for person detection + color/shape helmet logic",
            "Storage: uploads/ and results/",
        ],
        [
            "Endpoints:",
            "- GET /",
            "- POST /detect",
            "- POST /detect_video",
            "- GET /realtime",
            "- GET /api/stats",
        ],
    )

    # AI Model Details
    add_content_slide(
        prs,
        "AI Model Details",
        [
            "YOLOv8n for person detection (Ultralytics)",
            "Helmet detection via HSV color + Hough Circle",
            "Combined score = 0.7*color + 0.3*shape",
            "Configurable thresholds and color ranges",
        ],
        [
            "Tech Stack:",
            "- Python, Flask",
            "- OpenCV, NumPy",
            "- PyTorch, Ultralytics",
            "- HTML, CSS, JavaScript",
        ],
    )

    # Demo Flow
    add_bulleted_slide(
        prs,
        "Demo Flow",
        [
            "Upload image/video or open real-time page",
            "Backend saves file to uploads/",
            "Detector runs person + helmet analysis",
            "Annotated outputs saved in results/",
            "Dashboard displays detections & stats",
        ],
    )

    # Screens / Images if available
    for path in ["dd.jpeg", "sss.jpg", "sm.png"]:
        if os.path.exists(path):
            add_image_slide(prs, f"Sample: {path}", path)

    # API Endpoints
    add_content_slide(
        prs,
        "API Endpoints",
        [
            "POST /detect: image detection (multipart/form-data)",
            "POST /detect_video: video analysis",
            "GET /api/stats: statistics JSON",
            "GET /results/<file>: result media",
        ],
        [
            "Response (detect):",
            "- success, detections, summary",
            "- result_filename",
            "Summary: counts + compliance",
        ],
    )

    # Results & Analytics
    add_bulleted_slide(
        prs,
        "Results & Analytics",
        [
            "Per-image summary: persons, helmet/no-helmet, compliance %",
            "Video summary: frames analyzed, average compliance",
            "Live stats on realtime page",
            "Export JSON stats via /api/stats",
        ],
    )

    # Setup & Run
    add_content_slide(
        prs,
        "Setup & Run",
        [
            "pip install -r requirements.txt",
            "python run.py (auto opens browser)",
            "Dashboard: http://localhost:5000",
            "Test: python demo.py, python test_system.py",
        ],
        [
            "Folders (auto):",
            "- uploads/",
            "- results/",
            "Edit detection thresholds in helmet_detector.py",
        ],
    )

    # Future Work
    add_bulleted_slide(
        prs,
        "Future Enhancements",
        [
            "Train a dedicated helmet classifier",
            "Edge deployment & GPU acceleration",
            "Multi-camera aggregation & alerts",
            "Role-based access & audit trails",
        ],
    )

    # Thank you
    add_bulleted_slide(
        prs,
        "Thank You",
        [
            "Q&A",
            "Repo structure documented in README.md",
            "Contact: Project Owner",
        ],
    )

    prs.save(OUTPUT_FILE)
    print(f"✅ Presentation generated: {OUTPUT_FILE}")


if __name__ == "__main__":
    build_presentation()

